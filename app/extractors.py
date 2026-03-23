from __future__ import annotations

import importlib
import os
import re
import shutil
import subprocess
import tempfile
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable

from docx import Document
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from PIL import Image

from .document_compare.extraction_layout import ExtractionBlock, ExtractionLayout

ALWAYS_SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
LEGACY_OFFICE_EXTENSIONS = {".doc", ".xls", ".rtf"}
SUPPORTED_TEXT_EXTENSIONS = ALWAYS_SUPPORTED_EXTENSIONS | LEGACY_OFFICE_EXTENSIONS | IMAGE_EXTENSIONS
SUPPORTED_ENGINES = {"auto", "builtin", "docling"}
_LIST_PREFIX_RX = re.compile(r"^(?:[\-–—•▪◦·●○■□]+|(?:\(?\d+(?:\.\d+)*[\)\.]|[a-z]\)|[ivxlcdm]+[\)\.]))\s+", re.IGNORECASE)


class UnsupportedFormatError(ValueError):
    pass


class UnsupportedEngineError(ValueError):
    pass


class MissingDependencyError(RuntimeError):
    pass


@dataclass(slots=True)
class ExtractionResult:
    text: str
    engine: str
    quality_score: float
    metadata: dict[str, Any] = field(default_factory=dict)
    blocks: list[ExtractionBlock] = field(default_factory=list)
    quality_signals: dict[str, Any] = field(default_factory=dict)

    def to_quality_dict(self) -> dict[str, Any]:
        return {
            "engine": self.engine,
            "score": self.quality_score,
            "signals": dict(self.quality_signals),
            "metadata": {
                **dict(self.metadata),
                "block_count": len(self.blocks),
                "block_types": dict(Counter(block.block_type for block in self.blocks)),
            },
        }


def normalize_requested_engine(engine: str | None, *, strict: bool = False) -> str:
    normalized = str(engine or "auto").strip().lower()
    if normalized in SUPPORTED_ENGINES:
        return normalized
    if strict:
        raise UnsupportedEngineError(
            f"Engine no soportado: {normalized or 'vacío'}. Motores admitidos: {', '.join(sorted(SUPPORTED_ENGINES))}."
        )
    return "auto"


def docling_available() -> bool:
    return importlib.util.find_spec("docling") is not None


def detect_ocr_backend() -> str | None:
    if importlib.util.find_spec("rapidocr_onnxruntime") is not None:
        return "rapidocr"
    if importlib.util.find_spec("pytesseract") is not None and shutil.which("tesseract"):
        return "tesseract"
    return None


def pptx_available() -> bool:
    return importlib.util.find_spec("pptx") is not None


def get_pipeline_capabilities(*, soffice_path: str | None = None) -> dict[str, Any]:
    resolved_soffice = validate_soffice_option(soffice_path)
    ocr_backend = detect_ocr_backend()
    allowed_extensions = sorted(ext for ext in ALWAYS_SUPPORTED_EXTENSIONS if ext != ".pptx" or pptx_available())
    if resolved_soffice:
        allowed_extensions.extend(sorted(LEGACY_OFFICE_EXTENSIONS))
    if ocr_backend:
        allowed_extensions.extend(sorted(IMAGE_EXTENSIONS))

    conditional_support = {
        ext: {
            "requires": "soffice",
            "available": bool(resolved_soffice),
            "detail": "Requiere LibreOffice headless para convertir a un formato OOXML antes de extraer.",
        }
        for ext in sorted(LEGACY_OFFICE_EXTENSIONS)
    }
    conditional_support[".pptx"] = {
        "requires": "python-pptx",
        "available": pptx_available(),
        "detail": "Requiere la dependencia python-pptx para extraer texto de presentaciones .pptx.",
    }

    return {
        "allowed_extensions": allowed_extensions,
        "conditional_extensions": conditional_support,
        "engines": {
            "default": "auto",
            "available": ["auto", "builtin", *(["docling"] if docling_available() else [])],
            "declared": sorted(SUPPORTED_ENGINES),
            "availability": {
                "auto": True,
                "builtin": True,
                "docling": docling_available(),
            },
        },
        "soffice": {
            "available": bool(resolved_soffice),
            "path": resolved_soffice,
        },
        "ocr": {
            "available": bool(ocr_backend),
            "backend": ocr_backend,
        },
    }


def validate_extraction_request(path: str | Path, *, soffice_path: str | None = None, engine: str = "auto") -> tuple[str, str | None]:
    normalized_engine = normalize_requested_engine(engine, strict=True)
    ext = Path(path).suffix.lower()
    if ext not in SUPPORTED_TEXT_EXTENSIONS:
        raise UnsupportedFormatError(f"Formato no soportado: {ext or 'sin extensión'}")

    resolved_soffice = validate_soffice_option(soffice_path)
    if ext in LEGACY_OFFICE_EXTENSIONS and not resolved_soffice:
        raise MissingDependencyError(
            f"El formato {ext} solo está soportado condicionalmente y requiere LibreOffice/soffice disponible o una ruta válida en 'soffice'."
        )

    if ext in IMAGE_EXTENSIONS and not detect_ocr_backend():
        raise MissingDependencyError(
            f"El formato {ext} requiere un backend OCR disponible (RapidOCR o Tesseract)."
        )

    if ext == ".pptx" and not pptx_available():
        raise MissingDependencyError(
            "El formato .pptx requiere la dependencia python-pptx instalada en el backend."
        )

    if normalized_engine == "docling" and not docling_available():
        raise MissingDependencyError(
            "Se solicitó engine=docling, pero Docling no está disponible en este despliegue."
        )

    if normalized_engine == "docling" and ext != ".pdf":
        raise UnsupportedEngineError(
            f"El engine docling solo está disponible para PDF en este pipeline. Formato recibido: {ext or 'sin extensión'}."
        )

    return normalized_engine, resolved_soffice


def validate_soffice_option(value: str | None) -> str | None:
    for candidate in [value, os.getenv("SOFFICE_PATH"), shutil.which("soffice"), shutil.which("libreoffice")]:
        if not candidate:
            continue
        candidate = str(candidate).strip().strip('"').strip("'")
        if not candidate:
            continue
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
        if Path(candidate).exists():
            return str(Path(candidate))
    return None


def extraction_is_reliable(opts: dict[str, Any] | None) -> bool:
    qa = (opts or {}).get("_qualityA") or {}
    qb = (opts or {}).get("_qualityB") or {}
    return min(float(qa.get("score") or 0.0), float(qb.get("score") or 0.0)) >= 0.5


def _decode_text_bytes(content: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def extract_document_text(path: str, *, soffice_path: str | None = None, drop_headers: bool = True, engine: str = "auto") -> tuple[str, dict[str, Any]]:
    result = extract_document_result(path, soffice_path=soffice_path, drop_headers=drop_headers, engine=engine)
    return result.text, result.to_quality_dict()


def extract_document_result(path: str, *, soffice_path: str | None = None, drop_headers: bool = True, engine: str = "auto") -> ExtractionResult:
    source = Path(path)
    ext = source.suffix.lower()
    requested_engine, resolved_soffice = validate_extraction_request(source, soffice_path=soffice_path, engine=engine)

    if ext == ".txt":
        result = _extract_txt(source)
    elif ext == ".docx":
        result = _extract_docx(source)
    elif ext == ".xlsx":
        result = _extract_xlsx(source)
    elif ext == ".pptx":
        result = _extract_pptx(source)
    elif ext == ".pdf":
        result = _extract_pdf(source, drop_headers=drop_headers, engine=requested_engine)
    elif ext in IMAGE_EXTENSIONS:
        result = _extract_image(source)
    elif ext in LEGACY_OFFICE_EXTENSIONS:
        if not resolved_soffice:
            raise MissingDependencyError(f"Se requiere LibreOffice/soffice para extraer {ext}")
        conversion_result = _convert_legacy(source, resolved_soffice)
        cleanup_dir: Path | None = None
        if isinstance(conversion_result, tuple):
            converted, cleanup_dir = conversion_result
        else:
            converted = conversion_result
        try:
            nested = extract_document_result(
                str(converted),
                soffice_path=resolved_soffice,
                drop_headers=drop_headers,
                engine=requested_engine,
            )
        finally:
            try:
                if cleanup_dir is not None and cleanup_dir.exists():
                    shutil.rmtree(cleanup_dir, ignore_errors=True)
            except Exception:
                pass
        metadata = {
            **nested.metadata,
            "source_format": nested.metadata.get("source_format") or converted.suffix.lower().lstrip("."),
            "source_format_real": ext.lstrip("."),
            "conversion": {
                "applied": True,
                "converter": "libreoffice",
                "soffice_path": resolved_soffice,
                "from_extension": ext,
                "to_extension": converted.suffix.lower(),
            },
            "engine_requested": requested_engine,
            "engine_used": nested.metadata.get("engine_used", nested.engine),
        }
        return ExtractionResult(
            text=nested.text,
            engine=nested.engine,
            quality_score=nested.quality_score,
            metadata=metadata,
            blocks=nested.blocks,
            quality_signals=nested.quality_signals,
        )
    else:
        raise UnsupportedFormatError(f"No hay extractor para {ext or 'sin extensión'}")

    metadata = {
        **result.metadata,
        "source_format": result.metadata.get("source_format") or ext.lstrip("."),
        "source_format_real": ext.lstrip("."),
        "conversion": result.metadata.get("conversion") or {"applied": False},
        "engine_requested": requested_engine,
        "engine_used": result.metadata.get("engine_used", result.engine),
    }
    return ExtractionResult(
        text=result.text,
        engine=result.engine,
        quality_score=result.quality_score,
        metadata=metadata,
        blocks=result.blocks,
        quality_signals=result.quality_signals,
    )


def _extract_txt(path: Path) -> ExtractionResult:
    content = path.read_bytes()
    text = _decode_text_bytes(content)
    blocks = _simple_blocks_from_text(text, source_engine="txt")
    return _result_from_layout(
        source_engine="txt",
        extraction_engine="builtin",
        blocks=blocks,
        quality_score=1.0,
        metadata={"source_format": "txt", "engine_used": "builtin"},
    )


def _extract_docx(path: Path) -> ExtractionResult:
    document = Document(str(path))
    blocks: list[ExtractionBlock] = []

    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            para_blocks = _docx_blocks_from_paragraph(block)
            blocks.extend(para_blocks)
            if _paragraph_has_section_break(block):
                blocks.append(ExtractionBlock(text="", page=None, block_type="section_break", source_engine="docx"))
        elif isinstance(block, Table):
            for row_index, row in enumerate(block.rows, start=1):
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if not cells:
                    continue
                blocks.append(
                    ExtractionBlock(
                        text=" | ".join(cells),
                        page=None,
                        block_type="table_row",
                        source_engine="docx",
                        metadata={"row_index": row_index, "column_count": len(cells)},
                    )
                )

    metadata = {
        "source_format": "docx",
        "paragraph_count": sum(1 for block in blocks if block.block_type == "paragraph"),
        "list_item_count": sum(1 for block in blocks if block.block_type == "list_item"),
        "table_row_count": sum(1 for block in blocks if block.block_type == "table_row"),
        "section_break_count": sum(1 for block in blocks if block.block_type == "section_break"),
    }
    return _result_from_layout(source_engine="docx", extraction_engine="builtin", blocks=blocks, quality_score=0.98, metadata={**metadata, "engine_used": "builtin"})


def _extract_xlsx(path: Path) -> ExtractionResult:
    workbook = load_workbook(str(path), data_only=True)
    blocks: list[ExtractionBlock] = []
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        blocks.append(ExtractionBlock(text=sheet.title, page=sheet_index, block_type="header", source_engine="xlsx", metadata={"sheet": sheet.title}))
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                blocks.append(
                    ExtractionBlock(
                        text=" | ".join(values),
                        page=sheet_index,
                        block_type="table_row",
                        source_engine="xlsx",
                        metadata={"sheet": sheet.title, "row_index": row_index, "column_count": len(values)},
                    )
                )
    return _result_from_layout(source_engine="xlsx", extraction_engine="builtin", blocks=blocks, quality_score=0.96, metadata={"source_format": "xlsx", "sheet_count": len(workbook.worksheets), "engine_used": "builtin"})


def _extract_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation

    presentation = Presentation(str(path))
    blocks: list[ExtractionBlock] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        blocks.append(ExtractionBlock(text=f"Slide {slide_index}", page=slide_index, block_type="page", source_engine="pptx", metadata={"drop_in_canonical": True}))
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = getattr(shape, "text", "")
            if text and text.strip():
                for line in [line.strip() for line in text.splitlines() if line.strip()]:
                    blocks.append(
                        ExtractionBlock(
                            text=line,
                            page=slide_index,
                            block_type="paragraph",
                            source_engine="pptx",
                            metadata={"shape_index": shape_index},
                        )
                    )
    return _result_from_layout(source_engine="pptx", blocks=blocks, quality_score=0.95, metadata={"slide_count": len(presentation.slides)})


def _extract_pdf(path: Path, *, drop_headers: bool = True, engine: str = "auto") -> ExtractionResult:
    selected_engine = "docling" if engine == "docling" else "builtin"
    if selected_engine == "docling":
        return _extract_pdf_with_docling(path, drop_headers=drop_headers)
    return _extract_pdf_with_builtin(path, drop_headers=drop_headers)


def _extract_pdf_with_builtin(path: Path, *, drop_headers: bool = True) -> ExtractionResult:
    fitz = _load_pymupdf()
    doc = fitz.open(str(path))
    try:
        blocks: list[ExtractionBlock] = []
        page_heights: dict[int, float] = {}
        for page_index, page in enumerate(doc, start=1):
            page_heights[page_index] = float(page.rect.height or 0.0)
            blocks.append(ExtractionBlock(text=f"Page {page_index}", page=page_index, block_type="page", source_engine="pdf", metadata={"drop_in_canonical": True}))
            page_dict = page.get_text("dict")
            blocks.extend(_pdf_blocks_from_page(page_dict, page_number=page_index, page_height=page_heights[page_index]))

        repeated = _mark_repeated_pdf_headers_and_footers(blocks, page_heights=page_heights, drop_headers=drop_headers)
        metadata = {
            "source_format": "pdf",
            "page_count": len(doc),
            "engine_used": "builtin",
            **repeated,
        }
        return _result_from_layout(source_engine="pdf", extraction_engine="builtin", blocks=blocks, quality_score=0.93, metadata=metadata)
    finally:
        doc.close()


def _extract_pdf_with_docling(path: Path, *, drop_headers: bool = True) -> ExtractionResult:
    if not docling_available():
        raise MissingDependencyError("Docling no está disponible para procesar PDFs con engine=docling.")

    from docling.document_converter import DocumentConverter  # type: ignore

    converter = DocumentConverter()
    converted = converter.convert(str(path))
    document = getattr(converted, "document", converted)
    markdown = ""
    if hasattr(document, "export_to_markdown"):
        markdown = str(document.export_to_markdown() or "")
    elif hasattr(document, "text"):
        markdown = str(getattr(document, "text") or "")
    else:
        markdown = str(document or "")
    blocks = _simple_blocks_from_text(markdown, source_engine="pdf")
    return _result_from_layout(
        source_engine="pdf",
        extraction_engine="docling",
        blocks=blocks,
        quality_score=0.94,
        metadata={
            "source_format": "pdf",
            "drop_headers": drop_headers,
            "engine_used": "docling",
            "conversion": {"applied": False},
        },
    )


def _extract_image(path: Path) -> ExtractionResult:
    if importlib.util.find_spec("rapidocr_onnxruntime") is None:
        return _extract_image_with_tesseract(path)

    from rapidocr_onnxruntime import RapidOCR

    engine = RapidOCR()
    result, _ = engine(str(path))
    blocks: list[ExtractionBlock] = []
    for index, item in enumerate(result or [], start=1):
        if len(item) < 2 or not item[1]:
            continue
        bbox = None
        if item[0] and len(item[0]) >= 4:
            xs = [point[0] for point in item[0]]
            ys = [point[1] for point in item[0]]
            bbox = (float(min(xs)), float(min(ys)), float(max(xs)), float(max(ys)))
        blocks.append(ExtractionBlock(text=str(item[1]).strip(), page=1, block_type="line", source_engine="rapidocr", bbox=bbox, metadata={"ocr_index": index}))
    return _result_from_layout(source_engine="rapidocr", extraction_engine="rapidocr", blocks=blocks, quality_score=0.75 if blocks else 0.0, metadata={"source_format": "image", "engine_used": "rapidocr"})


def _extract_image_with_tesseract(path: Path) -> ExtractionResult:
    if importlib.util.find_spec("pytesseract") is None:
        raise MissingDependencyError(
            "No hay un backend OCR disponible para imágenes (RapidOCR o Tesseract)."
        )

    import pytesseract  # type: ignore

    image = Image.open(path)
    try:
        text = pytesseract.image_to_string(image, lang="spa+eng").strip()
    except Exception:
        text = pytesseract.image_to_string(image).strip()
    blocks = _simple_blocks_from_text(text, source_engine="tesseract")
    return _result_from_layout(source_engine="tesseract", extraction_engine="tesseract", blocks=blocks, quality_score=0.68 if text else 0.0, metadata={"source_format": "image", "size": image.size, "engine_used": "tesseract"})


def _convert_legacy(path: Path, soffice_path: str) -> tuple[Path, Path]:
    target_dir = Path(tempfile.mkdtemp(prefix="comp_docs_convert_"))
    target_ext = ".docx" if path.suffix.lower() in {".doc", ".rtf"} else ".xlsx" if path.suffix.lower() == ".xls" else ".pptx"
    subprocess.run(
        [soffice_path, "--headless", "--convert-to", target_ext.lstrip('.'), "--outdir", str(target_dir), str(path)],
        check=True,
        capture_output=True,
        text=True,
    )
    matches = list(target_dir.glob(f"{path.stem}*{target_ext}"))
    if not matches:
        raise RuntimeError(f"No se pudo convertir {path.name} con LibreOffice")
    return matches[0], target_dir


def _load_pymupdf():
    pymupdf_spec = importlib.util.find_spec("pymupdf")
    if pymupdf_spec is not None:
        fitz = importlib.import_module("pymupdf")  # type: ignore
        if hasattr(fitz, "open"):
            return fitz

    fitz_spec = importlib.util.find_spec("fitz")
    if fitz_spec is not None:
        fitz = importlib.import_module("fitz")  # type: ignore
        if hasattr(fitz, "open"):
            return fitz

    raise RuntimeError(
        "No se pudo cargar PyMuPDF para procesar PDFs. "
        "Si aparece un error como `from frontend import *`, tienes instalado el paquete `fitz` incorrecto. "
        "Solución recomendada: `pip uninstall -y fitz` y después `pip install --upgrade pymupdf`."
    )


def _result_from_layout(*, source_engine: str, extraction_engine: str | None = None, blocks: list[ExtractionBlock], quality_score: float, metadata: dict[str, Any] | None = None) -> ExtractionResult:
    layout = ExtractionLayout(blocks=blocks, source_engine=source_engine, metadata=dict(metadata or {}))
    quality_signals = layout.quality_signals()
    metadata_dict = dict(metadata or {})
    metadata_dict.setdefault(
        "scan_like",
        source_engine == "rapidocr" or (
            source_engine == "pdf"
            and (
                float(quality_signals.get("avg_line_length") or 0.0) <= 22.0
                or float(quality_signals.get("layout_noise_score") or 0.0) >= 0.55
            )
        ),
    )
    metadata_dict.setdefault("source_engine", source_engine)
    metadata_dict.setdefault("canonical_reconstruction", "layout_blocks_v1")
    metadata_dict.setdefault("block_metadata_available", True)
    metadata_dict.setdefault(
        "block_provenance_hints",
        [
            "paragraph↔line-wrap",
            "table linearization",
            "header/footer repetition",
            "section or page breaks",
        ],
    )
    adjusted_score = _adjust_quality_score(base_score=quality_score, signals=quality_signals)
    return ExtractionResult(
        text=layout.canonical_text(),
        engine=extraction_engine or source_engine,
        quality_score=adjusted_score,
        metadata=metadata_dict,
        blocks=blocks,
        quality_signals=quality_signals,
    )


def _adjust_quality_score(*, base_score: float, signals: dict[str, Any]) -> float:
    penalty = 0.0
    penalty += min(0.18, float(signals.get("layout_noise_score") or 0.0) * 0.18)
    penalty += min(0.1, float(signals.get("table_like_density") or 0.0) * 0.1)
    if signals.get("has_repeated_headers"):
        penalty += 0.03
    return round(max(0.0, min(1.0, base_score - penalty)), 4)


def _simple_blocks_from_text(text: str, *, source_engine: str) -> list[ExtractionBlock]:
    blocks: list[ExtractionBlock] = []
    for chunk in [part.strip() for part in str(text or "").split("\n\n") if part.strip()]:
        block_type = "list_item" if _LIST_PREFIX_RX.match(chunk) else "paragraph"
        blocks.append(ExtractionBlock(text=chunk, page=None, block_type=block_type, source_engine=source_engine))
        for line in [line.strip() for line in chunk.splitlines() if line.strip()]:
            blocks.append(ExtractionBlock(text=line, page=None, block_type="line", source_engine=source_engine, metadata={"drop_in_canonical": True}))
    return blocks


def _iter_docx_blocks(document: Document) -> Iterable[Paragraph | Table]:
    body = document.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def _docx_blocks_from_paragraph(paragraph: Paragraph) -> list[ExtractionBlock]:
    text = paragraph.text.strip()
    blocks: list[ExtractionBlock] = []
    if _paragraph_has_page_break(paragraph):
        blocks.append(ExtractionBlock(text="", page=None, block_type="page_break", source_engine="docx"))
    if not text:
        return blocks

    style_name = str(getattr(paragraph.style, "name", "") or "").strip().lower()
    if style_name.startswith("heading") or style_name.startswith("title"):
        block_type = "header"
    elif _paragraph_is_list_item(paragraph, style_name=style_name):
        block_type = "list_item"
    else:
        block_type = "paragraph"

    blocks.append(
        ExtractionBlock(
            text=text,
            page=None,
            block_type=block_type,
            source_engine="docx",
            metadata={"style": style_name or None},
        )
    )
    for line in [line.strip() for line in text.splitlines() if line.strip()]:
        blocks.append(
            ExtractionBlock(
                text=line,
                page=None,
                block_type="line",
                source_engine="docx",
                metadata={"style": style_name or None, "drop_in_canonical": True},
            )
        )
    return blocks


def _paragraph_is_list_item(paragraph: Paragraph, *, style_name: str) -> bool:
    p_pr = paragraph._p.pPr
    has_num_pr = bool(p_pr is not None and p_pr.numPr is not None)
    return has_num_pr or any(token in style_name for token in ("list", "bullet", "number")) or bool(_LIST_PREFIX_RX.match(paragraph.text.strip()))


def _paragraph_has_page_break(paragraph: Paragraph) -> bool:
    return bool(paragraph._element.xpath(".//*[local-name()='br' and (@*[local-name()='type']='page')]") )


def _paragraph_has_section_break(paragraph: Paragraph) -> bool:
    return bool(paragraph._element.xpath("./*[local-name()='pPr']/*[local-name()='sectPr']"))


def _pdf_blocks_from_page(page_dict: dict[str, Any], *, page_number: int, page_height: float) -> list[ExtractionBlock]:
    blocks: list[ExtractionBlock] = []
    for raw_block in page_dict.get("blocks") or []:
        if raw_block.get("type") != 0:
            continue
        line_entries = _pdf_line_entries(raw_block)
        if not line_entries:
            continue
        paragraph_text = " ".join(entry["text"] for entry in line_entries).strip()
        paragraph_bbox = _bbox_from_items([entry["bbox"] for entry in line_entries if entry.get("bbox")])
        paragraph_type = _classify_pdf_group(line_entries, page_height=page_height)
        blocks.append(
            ExtractionBlock(
                text=paragraph_text,
                page=page_number,
                block_type=paragraph_type,
                source_engine="pdf",
                bbox=paragraph_bbox,
                metadata={"line_count": len(line_entries)},
            )
        )
        for line_index, entry in enumerate(line_entries, start=1):
            line_type = entry["block_type"]
            blocks.append(
                ExtractionBlock(
                    text=entry["text"],
                    page=page_number,
                    block_type=line_type,
                    source_engine="pdf",
                    bbox=entry.get("bbox"),
                    metadata={"line_index": line_index, "drop_in_canonical": True, "parent_block_type": paragraph_type},
                )
            )
    return blocks


def _pdf_line_entries(raw_block: dict[str, Any]) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    for raw_line in raw_block.get("lines") or []:
        spans = [span for span in raw_line.get("spans") or [] if str(span.get("text") or "").strip()]
        if not spans:
            continue
        text = " ".join(str(span.get("text") or "").strip() for span in spans if str(span.get("text") or "").strip()).strip()
        if not text:
            continue
        bbox = _bbox_from_items([tuple(span.get("bbox")) for span in spans if span.get("bbox")])
        entries.append(
            {
                "text": text,
                "bbox": bbox,
                "span_count": len(spans),
                "block_type": "table_row" if _looks_like_table_row(text=text, spans=spans) else ("list_item" if _LIST_PREFIX_RX.match(text) else "line"),
            }
        )
    return entries


def _classify_pdf_group(line_entries: list[dict[str, Any]], *, page_height: float) -> str:
    if all(entry["block_type"] == "table_row" for entry in line_entries):
        return "table_row"
    if line_entries and line_entries[0]["block_type"] == "list_item":
        return "list_item"
    bbox = _bbox_from_items([entry["bbox"] for entry in line_entries if entry.get("bbox")])
    if bbox and page_height:
        top_ratio = bbox[1] / page_height
        bottom_ratio = bbox[3] / page_height
        if top_ratio <= 0.1 and len(" ".join(entry["text"] for entry in line_entries).split()) <= 18:
            return "header"
        if bottom_ratio >= 0.9 and len(" ".join(entry["text"] for entry in line_entries).split()) <= 18:
            return "footer"
    return "paragraph"


def _looks_like_table_row(*, text: str, spans: list[dict[str, Any]]) -> bool:
    if text.count("|") >= 2 or text.count("\t") >= 2:
        return True
    if len(spans) < 3:
        return False
    xs = [float(span.get("bbox", [0, 0, 0, 0])[0]) for span in spans if span.get("bbox")]
    if len(xs) < 3:
        return False
    gaps = [curr - prev for prev, curr in zip(xs, xs[1:])]
    return sum(1 for gap in gaps if gap >= 24.0) >= 2


def _mark_repeated_pdf_headers_and_footers(blocks: list[ExtractionBlock], *, page_heights: dict[int, float], drop_headers: bool) -> dict[str, Any]:
    grouped: dict[str, list[ExtractionBlock]] = {}
    for block in blocks:
        if block.page is None or block.block_type not in {"paragraph", "header", "footer", "line", "list_item", "table_row"}:
            continue
        key = _repetition_key(block.text)
        if not key:
            continue
        grouped.setdefault(key, []).append(block)

    repeated_header_keys: set[str] = set()
    repeated_footer_keys: set[str] = set()
    for key, group in grouped.items():
        pages = {block.page for block in group if block.page is not None}
        if len(pages) < 2:
            continue
        top_hits = 0
        bottom_hits = 0
        for block in group:
            height = page_heights.get(block.page or 0) or 0.0
            if not block.bbox or not height:
                continue
            if (block.bbox[1] / height) <= 0.12:
                top_hits += 1
            if (block.bbox[3] / height) >= 0.88:
                bottom_hits += 1
        if top_hits >= 2:
            repeated_header_keys.add(key)
        if bottom_hits >= 2:
            repeated_footer_keys.add(key)

    for block in blocks:
        key = _repetition_key(block.text)
        if key in repeated_header_keys:
            block.block_type = "header"
            if drop_headers:
                block.metadata["drop_in_canonical"] = True
                block.metadata["is_repeated_header"] = True
        elif key in repeated_footer_keys:
            block.block_type = "footer"
            if drop_headers:
                block.metadata["drop_in_canonical"] = True
                block.metadata["is_repeated_footer"] = True

    return {
        "repeated_header_keys": len(repeated_header_keys),
        "repeated_footer_keys": len(repeated_footer_keys),
    }


def _bbox_from_items(items: Iterable[tuple[float, float, float, float] | list[float] | None]) -> tuple[float, float, float, float] | None:
    normalized = [tuple(float(value) for value in item) for item in items if item]
    if not normalized:
        return None
    xs0 = [item[0] for item in normalized]
    ys0 = [item[1] for item in normalized]
    xs1 = [item[2] for item in normalized]
    ys1 = [item[3] for item in normalized]
    return (min(xs0), min(ys0), max(xs1), max(ys1))


def _repetition_key(text: str) -> str:
    normalized = " ".join(str(text or "").casefold().split())
    normalized = re.sub(r"\d+", "#", normalized)
    return normalized if len(normalized) >= 3 else ""